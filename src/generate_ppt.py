import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_bullet(text_frame, text, level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    return p

prs = Presentation()
IMG_DIR = r"C:\Users\DAVID PC\Documents\credit-risk-scoring\reports\figures"

# --- Slide 1: Titre ---
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Modélisation du Risque de Crédit\n(Credit Risk Scoring)"
subtitle.text = "Projet PAEI\n\nPrésenté par : David Janvion HAMAYADJI NGOMNA\nEncadrant : Ing. Archange KOUMBA MOUITY"

# --- Slide 2: Contexte & Objectifs ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "1. Contexte & Objectifs du Projet"
tf = slide.placeholders[1].text_frame
tf.text = "Problématique :"
add_bullet(tf, "Comment anticiper le défaut de paiement d'un client bancaire pour sécuriser les octrois de crédits ?", 1)
add_bullet(tf, "Objectifs du projet :", 0)
add_bullet(tf, "Traiter et analyser 250 000 profils clients (Dataset Kaggle).", 1)
add_bullet(tf, "Construire un Indice de Risque fiable et interprétable.", 1)
add_bullet(tf, "Déployer une interface visuelle d'aide à la décision (Outil Back-Office).", 1)

# --- Slide 3: Le Pipeline de Données ---
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "2. Préparation et Feature Engineering"
tf = slide.placeholders[1].text_frame
tf.text = "Transformation de la donnée brute en donnée intelligente :"
add_bullet(tf, "Création de variables expertes : Calcul du ratio Dette/Revenu (LoanToIncomeRatio).", 1)
add_bullet(tf, "Encodage des données (One-Hot Encoding) :", 1)
add_bullet(tf, "Transformation des textes (Marié, Célibataire) en variables binaires mathématiques (0 ou 1).", 2)
add_bullet(tf, "Mise à l'échelle (StandardScaler) :", 1)
add_bullet(tf, "Normalisation des montants pour éviter que les hauts revenus n'écrasent les petits taux d'intérêt.", 2)

# --- Slide 4: Analyse Bivariée ---
slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
slide.shapes.title.text = "3. Exploration : Analyse Bivariée (Preuve Visuelle)"
img_path = os.path.join(IMG_DIR, "bivar_numeriques.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.2), Inches(1.5), width=Inches(9.5))
    
# --- Slide 5: Matrice de Corrélation ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "4. Mesure d'intensité : Matrice de Corrélation"
img_path = os.path.join(IMG_DIR, "correlation_matrix.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.5), height=Inches(5.5))

# --- Slide 6: Choix Algorithmiques ---
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "5. Justification des Algorithmes (IA)"
tf = slide.placeholders[1].text_frame
tf.text = "Deux algorithmes utilisés avec des objectifs distincts :"
add_bullet(tf, "Le Random Forest (Forêt Aléatoire) :", 0)
add_bullet(tf, "Utilisé en phase de recherche car robuste aux valeurs extrêmes.", 1)
add_bullet(tf, "Permet d'extraire mathématiquement l'Importance des Variables (Feature Importance).", 1)
add_bullet(tf, "La Régression Logistique (Production) :", 0)
add_bullet(tf, "Utilisée pour le livrable final car elle offre une Transparence Totale (exigée en banque).", 1)
add_bullet(tf, "Génère des poids exacts exportables directement vers une application Web.", 1)

# --- Slide 7: Importance des Variables ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "6. Les déclencheurs de défaut (Random Forest)"
img_path = os.path.join(IMG_DIR, "feature_importance_rf.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(5.5))

# --- Slide 8: Indice de Risque ML ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "7. L'Indice de Risque Machine Learning"
img_path = os.path.join(IMG_DIR, "risk_index_ml.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# --- Slide 9: Évaluation des Performances ---
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "8. Évaluation et Score AUC-ROC"
tf = slide.placeholders[1].text_frame
tf.text = "Méthodologie d'évaluation (Test sur données inconnues) :"
add_bullet(tf, "Séparation : 80% des clients pour l'apprentissage, 20% pour le test de validation.", 1)
add_bullet(tf, "Performance obtenue : AUC = 0.76", 0)
add_bullet(tf, "Interprétation métier : Si on prend au hasard un 'Mauvais payeur' et un 'Bon payeur', le modèle classera correctement le profil risqué dans 76% des cas.", 1)
add_bullet(tf, "C'est un score robuste et suffisant pour automatiser un premier filtrage bancaire.", 1)

# --- Slide 10: Le Livrable Final ---
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "9. Mise en Production (Simulateur Web)"
tf = slide.placeholders[1].text_frame
tf.text = "Un déploiement moderne et Serverless :"
add_bullet(tf, "Le script Python extrait l'équation de l'IA (Poids) au format JSON.", 1)
add_bullet(tf, "Création d'un Dashboard interactif en HTML / JS / TailwindCSS.", 1)
add_bullet(tf, "Aide à la décision : L'interface affiche une jauge de risque instantanée et un code couleur (Vert, Orange, Rouge).", 1)
add_bullet(tf, "Démonstration du simulateur...", 0)

# Sauvegarde
output_path = r"C:\Users\DAVID PC\Documents\credit-risk-scoring\src\Presentation_Soutenance_PAEI.pptx"
prs.save(output_path)
print(f"PowerPoint généré avec succès : {output_path}")
